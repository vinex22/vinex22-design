from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
IMG = ROOT / "assets" / "images"
OUT = ROOT / "output" / "naturz-exhibition-clean-editable.pptx"
NATIVE = ROOT / "editable" / "native-assets"
NATIVE.mkdir(parents=True, exist_ok=True)

COLORS = {
    "green": RGBColor(21, 134, 74),
    "deep": RGBColor(12, 75, 46),
    "gold": RGBColor(243, 179, 47),
    "cream": RGBColor(255, 247, 230),
    "white": RGBColor(255, 253, 245),
    "ink": RGBColor(27, 33, 27),
    "muted": RGBColor(92, 106, 92),
    "rule": RGBColor(211, 201, 170),
}

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

ASSETS = {
    "logo": "logo.webp",
    "hero": "hero-products.webp",
    "oil_jerry": "oil-jerry.webp",
    "oil_tin": "oil-tin.webp",
    "oil_bottle": "oil-bottle.webp",
    "shortening": "shortening.webp",
    "ghee": "ghee.webp",
    "glycerin": "glycerin.webp",
    "chocolate": "chocolate-fat.webp",
    "frying": "frying-fat.webp",
    "stearin": "palm-stearin.webp",
    "factory": "production-line.webp",
    "iso": "iso.webp",
    "gmp": "gmp.webp",
    "mpob": "mpob.webp",
    "ms": "ms-standard.webp",
}


def convert_asset(name: str) -> Path:
    src = IMG / ASSETS[name]
    dst = NATIVE / f"{name}.png"
    if not dst.exists() or dst.stat().st_mtime < src.stat().st_mtime:
        with Image.open(src) as im:
            im.convert("RGBA").save(dst)
    return dst


for key in ASSETS:
    convert_asset(key)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=24, color=None, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_title(slide, text, x=0.7, y=1.35, w=6.2, h=1.8, size=42, color=None):
    return add_text(slide, text, x, y, w, h, size=size, color=color or COLORS["ink"],
                    bold=True, font="Georgia", line_spacing=0.92)


def add_body(slide, text, x=0.72, y=3.15, w=5.3, h=0.9, size=18, color=None):
    return add_text(slide, text, x, y, w, h, size=size, color=color or COLORS["muted"],
                    font="Aptos", line_spacing=1.05)


def add_kicker(slide, text, x=0.72, y=1.03, w=4.6, color=None):
    return add_text(slide, text.upper(), x, y, w, 0.25, size=9.5, color=color or COLORS["green"],
                    bold=True, font="Aptos")


def add_logo(slide, dark=False):
    pic = slide.shapes.add_picture(str(convert_asset("logo")), Inches(0.68), Inches(0.34), width=Inches(1.18))
    if dark:
        # Add a white backing so the official green logo remains readable.
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(0.29), Inches(1.32), Inches(0.62))
        bg.fill.solid(); bg.fill.fore_color.rgb = COLORS["cream"]
        bg.line.color.rgb = COLORS["cream"]
        slide.shapes._spTree.remove(pic._element)
        slide.shapes.add_picture(str(convert_asset("logo")), Inches(0.68), Inches(0.34), width=Inches(1.18))


def add_footer(slide, label, page, dark=False):
    color = COLORS["white"] if dark else COLORS["muted"]
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(6.86), Inches(11.98), Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(116, 147, 116) if dark else COLORS["rule"]
    line.line.color.rgb = line.fill.fore_color.rgb
    add_text(slide, label.upper(), 0.7, 6.94, 6, 0.2, size=8, color=color, bold=True)
    add_text(slide, f"{page:02d} / 10", 11.7, 6.94, 0.95, 0.2, size=8, color=color, bold=True, align=PP_ALIGN.RIGHT)


def add_picture_contain(slide, asset, x, y, w, h):
    path = convert_asset(asset)
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w
        height = w / img_ratio
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        left = x + (w - width) / 2
        top = y
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_panel(slide, x, y, w, h, fill="white", border=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = COLORS[fill]
    shp.line.color.rgb = COLORS["rule"] if border else COLORS[fill]
    shp.line.width = Pt(1)
    return shp


def add_card(slide, x, y, w, h, title, subtitle, asset):
    add_panel(slide, x, y, w, h)
    add_picture_contain(slide, asset, x + 0.38, y + 0.22, w - 0.76, 1.35)
    add_text(slide, title, x + 0.24, y + 1.72, w - 0.48, 0.35, size=17, color=COLORS["ink"], bold=True, font="Georgia", align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, x + 0.3, y + 2.12, w - 0.6, 0.35, size=10.8, color=COLORS["muted"], align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# 01
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_kicker(s, "Exhibition promotion")
add_title(s, "Malaysia's palm oil expertise, ready for your brand.", size=41)
add_body(s, "Cooking oil and palm derivative products for food manufacturers, distributors and OEM partners.", size=17.2)
add_picture_contain(s, "hero", 7.25, 1.12, 4.55, 3.75)
for x, label, val in [(0.72, "Founded", "2020"), (2.8, "Base", "Malaysia"), (4.85, "Contact", "info@naturzindustries.com")]:
    add_text(s, label.upper(), x, 5.62, 1.9, 0.18, size=8.5, color=COLORS["muted"], bold=True)
    add_text(s, val, x, 5.9, 3.0 if label == "Contact" else 1.5, 0.38, size=18 if label != "Contact" else 15, color=COLORS["ink"], bold=True, font="Georgia")
add_footer(s, "Naturz Industries Sdn Bhd", 1)

# 02
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_kicker(s, "Best for your plate")
add_title(s, "Naturz's best for your plate.", size=48)
add_body(s, "Dependable vegetable oil products for frying, baking, cooking and food processing.", size=18)
oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), Inches(1.45), Inches(2.85), Inches(2.85)); oval.fill.solid(); oval.fill.fore_color.rgb = COLORS["gold"]; oval.line.color.rgb = COLORS["gold"]
add_picture_contain(s, "oil_jerry", 7.4, 1.05, 4.0, 3.75)
add_text(s, "Reliable. Practical. Export-ready.", 7.28, 5.18, 4.2, 0.45, size=22, color=COLORS["deep"], bold=True, font="Georgia", align=PP_ALIGN.CENTER)
add_footer(s, "Brand promise", 2)

# 03
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_title(s, "One supplier. Multiple food applications.", 0.7, 0.95, 7.3, 0.72, size=31)
add_body(s, "Product lines for bulk buyers, distributors, food service, bakery, confectionery and industrial food use.", 0.72, 1.78, 7.2, 0.45, size=13.5)
add_card(s, 0.68, 2.62, 2.45, 2.58, "RBD Palm Olein", "CP5, CP6, CP7, CP8, CP10", "oil_tin")
add_card(s, 3.32, 2.62, 2.45, 2.58, "Shortening", "Bakery and frying needs", "shortening")
add_card(s, 5.96, 2.62, 2.45, 2.58, "Vegetable Ghee", "Cooking and formulation", "ghee")
add_card(s, 8.6, 2.62, 2.45, 2.58, "Glycerin", "Palm derivative category", "glycerin")
add_footer(s, "Product range", 3)

# 04
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_kicker(s, "Core product")
add_title(s, "RBD Palm Olein for cooking, frying and food processing.", size=37)
add_body(s, "Naturz lists a wide range of palm olein grades for OEM cooking oil service and export buyers.", size=16.5)
add_picture_contain(s, "oil_bottle", 7.55, 1.22, 1.75, 3.35)
add_picture_contain(s, "oil_jerry", 9.16, 1.0, 2.35, 3.72)
for i, grade in enumerate(["CP5", "CP6", "CP7", "CP8", "CP10"]):
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72 + i * 1.12), Inches(5.35), Inches(0.86), Inches(0.45))
    pill.fill.solid(); pill.fill.fore_color.rgb = COLORS["deep"]; pill.line.color.rgb = COLORS["deep"]
    add_text(s, grade, 0.72 + i * 1.12, 5.45, 0.86, 0.18, size=13, color=COLORS["white"], bold=True, font="Georgia", align=PP_ALIGN.CENTER)
add_footer(s, "RBD Palm Olein", 4)

# 05
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_title(s, "Specialty fats and derivatives for food applications.", 0.7, 0.95, 7.2, 0.85, size=31)
add_body(s, "A broader portfolio supports bakery, confectionery, food service and industrial formulation conversations.", 0.72, 1.85, 7.0, 0.45, size=13.5)
add_card(s, 1.05, 2.64, 2.7, 2.45, "Chocolate Fat", "Specialty fat category", "chocolate")
add_card(s, 5.05, 2.64, 2.7, 2.45, "Frying Fat", "Food-service use", "frying")
add_card(s, 9.05, 2.64, 2.7, 2.45, "Palm Stearin", "Palm derivative line", "stearin")
add_footer(s, "Specialty fats", 5)

# 06
s = prs.slides.add_slide(blank); set_bg(s, COLORS["deep"]); add_logo(s, dark=True)
add_kicker(s, "Manufacturing + export", color=COLORS["gold"])
add_title(s, "Built for B2B supply conversations.", size=41, color=COLORS["white"])
add_body(s, "Naturz describes its business as manufacturing, marketing and exporting palm oil derivatives.", size=17, color=RGBColor(231, 239, 228))
add_picture_contain(s, "factory", 7.15, 1.35, 4.55, 2.95)
for x, label, val in [(0.72, "Founded", "2020"), (2.85, "Base", "Malaysia"), (4.95, "Focus", "Palm derivatives")]:
    add_text(s, label.upper(), x, 5.62, 1.6, 0.18, size=8.5, color=RGBColor(200, 218, 200), bold=True)
    add_text(s, val, x, 5.9, 2.0, 0.36, size=18, color=COLORS["white"], bold=True, font="Georgia")
add_footer(s, "Production readiness", 6, dark=True)

# 07
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_kicker(s, "OEM + export")
add_title(s, "Move from product interest to quotation.", size=39)
add_body(s, "Capture grade, pack format, volume, destination and documentation needs at the booth.", size=17)
steps = [("1", "Choose product and grade", "RBD palm olein, shortening, ghee, frying fat or derivative category."), ("2", "Choose pack format", "Bottle, tin, jerry can, carton or buyer-specific packaging discussion."), ("3", "Start quotation", "Share market, volume, documentation needs and timeline.")]
for idx, (num, title, txt) in enumerate(steps):
    y = 1.35 + idx * 1.25
    add_panel(s, 7.15, y, 4.4, 0.9)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.38), Inches(y + 0.22), Inches(0.45), Inches(0.45)); circ.fill.solid(); circ.fill.fore_color.rgb = COLORS["deep"]; circ.line.color.rgb = COLORS["deep"]
    add_text(s, num, 7.38, y + 0.31, 0.45, 0.16, size=11, color=COLORS["white"], bold=True, font="Georgia", align=PP_ALIGN.CENTER)
    add_text(s, title, 8.02, y + 0.18, 3.1, 0.22, size=16, color=COLORS["ink"], bold=True, font="Georgia")
    add_text(s, txt, 8.02, y + 0.48, 3.05, 0.28, size=10.5, color=COLORS["muted"])
add_footer(s, "OEM + export", 7)

# 08
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_kicker(s, "Buyer confidence")
add_title(s, "Quality conversations start before quotation.", size=39)
add_body(s, "Use certification and standards cues to guide buyers toward the right documentation questions.", size=17)
for x, y, asset in [(7.1, 1.15, "iso"), (9.4, 1.15, "gmp"), (7.1, 3.0, "mpob"), (9.4, 3.0, "ms")]:
    add_panel(s, x, y, 1.85, 1.36)
    add_picture_contain(s, asset, x + 0.2, y + 0.18, 1.45, 1.0)
add_panel(s, 0.72, 5.05, 5.35, 0.55)
add_text(s, "Confirm certificate scope and validity per product, order and destination market.", 0.95, 5.23, 4.85, 0.18, size=10.5, color=COLORS["muted"], bold=True)
add_footer(s, "Quality cues", 8)

# 09
s = prs.slides.add_slide(blank); set_bg(s, COLORS["cream"]); add_logo(s)
add_kicker(s, "Visit our booth")
add_title(s, "Ask for grade, packaging and export quotation.", size=39)
add_body(s, "Bring product category, market, volume, packaging, documentation and delivery timeline.", size=17)
contacts = [("Email", "info@naturzindustries.com"), ("Phone", "+603 5569 1581"), ("Mobile", "+6011 1115 4642"), ("Malaysia office", "UOA Business Park, Shah Alam, Selangor")]
for idx, (label, val) in enumerate(contacts):
    y = 1.2 + idx * 1.02
    add_panel(s, 7.1, y, 4.45, 0.72)
    add_text(s, label.upper(), 7.35, y + 0.16, 3.5, 0.12, size=8.5, color=COLORS["green"], bold=True)
    add_text(s, val, 7.35, y + 0.38, 3.65, 0.18, size=14 if label != "Malaysia office" else 11.5, color=COLORS["ink"], bold=True, font="Georgia")
add_footer(s, "Contact", 9)

# 10
s = prs.slides.add_slide(blank); set_bg(s, COLORS["deep"]); add_logo(s, dark=True)
add_kicker(s, "Naturz Industries", color=COLORS["gold"])
add_title(s, "Let's build your next cooking oil supply program.", size=42, color=COLORS["white"])
add_body(s, "Talk to our team about product grade, packaging and export requirements.", size=18, color=RGBColor(231, 239, 228))
add_picture_contain(s, "hero", 7.0, 1.35, 4.6, 3.55)
add_text(s, "naturzindustries.com", 0.72, 5.35, 4.8, 0.35, size=22, color=COLORS["gold"], bold=True, font="Georgia")
add_text(s, "info@naturzindustries.com · +603 5569 1581", 0.72, 5.86, 5.8, 0.22, size=13, color=COLORS["white"], bold=True)
add_footer(s, "Closing CTA", 10, dark=True)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT}")
print(f"slides {len(prs.slides)}")
